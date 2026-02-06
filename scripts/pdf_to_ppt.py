#!/usr/bin/env python3
"""
PDF 转 PowerPoint 转换器
使用 pdfplumber 和 python-pptx 实现 PDF 到 PPTX 的转换
"""

from pathlib import Path
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re


def pdf_to_ppt(input_path: str, output_path: str) -> dict:
    """
    将 PDF 文件转换为 PowerPoint 演示文稿
    
    Args:
        input_path: PDF 文件路径
        output_path: 输出 PPT 文件路径
    
    Returns:
        dict: 转换结果信息
    """
    result = {
        "success": False,
        "pages": 0,
        "message": ""
    }
    
    try:
        # 创建 PowerPoint 演示文稿
        prs = Presentation()
        
        # 打开 PDF
        with pdfplumber.open(input_path) as pdf:
            result["pages"] = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                # 创建一个新幻灯片（空白布局）
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 添加页面标题
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
                )
                title_frame = title_box.text_frame
                title_frame.text = f"第 {page_num} 页 / 共 {len(pdf.pages)} 页"
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(14)
                title_para.font.color.rgb = RGBColor(100, 100, 100)
                title_para.alignment = PP_ALIGN.CENTER
                
                # 提取文本
                text = page.extract_text()
                
                if text:
                    # 清理文本
                    text = clean_text(text)
                    
                    # 创建文本框
                    text_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.0), Inches(9), Inches(5)
                    )
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    # 按段落分割
                    paragraphs = text.split('\n')
                    
                    for i, para in enumerate(paragraphs):
                        para = para.strip()
                        if not para:
                            continue
                        
                        # 处理标题（短行）
                        if len(para) < 50 and not para.endswith(('。', '！', '？', ')', ']', '.')):
                            if i == 0:
                                # 第一段作为大标题
                                p = text_frame.paragraphs[0]
                                p.text = para
                                p.font.size = Pt(24)
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(0, 51, 102)
                                p.space_before = Pt(12)
                            else:
                                # 添加新段落作为小标题
                                p = text_frame.add_paragraph()
                                p.text = para
                                p.font.size = Pt(18)
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(0, 102, 204)
                                p.space_before = Pt(18)
                        else:
                            # 普通段落
                            if i == 0 and text_frame.paragraphs[0].text == "":
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = para
                            p.font.size = Pt(16)
                            p.font.color.rgb = RGBColor(0, 0, 0)
                            p.space_before = Pt(6)
                
                # 提取表格（如果有）
                tables = page.extract_tables()
                if tables:
                    # 添加表格标题
                    table_title = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
                    )
                    tt_frame = table_title.text_frame
                    tt_frame.text = f"表格数据"
                    tt_frame.paragraphs[0].font.size = Pt(12)
                    tt_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                    
                    for table_idx, table in enumerate(tables):
                        if not table or not table[0]:
                            continue
                        
                        # 计算表格尺寸
                        rows = len(table)
                        cols = len(table[0])
                        
                        # 限制行列数（防止过大）
                        if rows > 50 or cols > 10:
                            continue
                        
                        # 添加表格
                        left = Inches(0.5)
                        top = Inches(5.5) + Inches(table_idx * 0.5)
                        width = Inches(9)
                        height = Inches(0.8)
                        
                        # 检查是否超出页面
                        if top + Inches(rows * 0.5) > Inches(7):
                            # 新建幻灯片
                            slide_layout = prs.slide_layouts[6]
                            slide = prs.slides.add_slide(slide_layout)
                            top = Inches(0.5)
                        
                        table_shape = slide.shapes.add_table(
                            rows, cols, left, top, width, height
                        )
                        table = table_shape.table
                        
                        # 设置列宽
                        for col_idx in range(cols):
                            table.columns[col_idx].width = Inches(9 / cols)
                        
                        # 填充数据
                        for row_idx, row in enumerate(table):
                            for col_idx, cell in enumerate(row):
                                if cell:
                                    cell.text = str(cell)
                                    # 第一行加粗（表头）
                                    if row_idx == 0:
                                        cell.text_frame.paragraphs[0].font.bold = True
                                        cell.text_frame.paragraphs[0].font.size = Pt(12)
                                    else:
                                        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # 保存演示文稿
        prs.save(output_path)
        
        result["success"] = True
        result["message"] = f"转换成功！共 {result['pages']} 页"
        
    except Exception as e:
        result["message"] = f"转换失败: {str(e)}"
        import traceback
        traceback.print_exc()
    
    return result


def clean_text(text: str) -> str:
    """清理提取的文本"""
    # 移除多余的空白字符
    text = re.sub(r'\s+', ' ', text)
    # 修复常见的 PDF 提取问题
    text = text.replace('\x00', '')
    # 移除特殊字符
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
    return text.strip()


def convert_batch(input_dir: str, output_dir: str) -> list:
    """
    批量转换 PDF 文件
    
    Args:
        input_dir: 输入文件夹路径
        output_dir: 输出文件夹路径
    
    Returns:
        list: 每个文件的转换结果列表
    """
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    
    results = []
    pdf_files = list(input_path.glob("*.pdf"))
    
    for pdf_file in pdf_files:
        output_file = output_path / f"{pdf_file.stem}.pptx"
        result = pdf_to_ppt(str(pdf_file), str(output_file))
        result["file"] = pdf_file.name
        results.append(result)
    
    return results


if __name__ == "__main__":
    # 测试转换
    import sys
    
    if len(sys.argv) > 1:
        input_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace('.pdf', '.pptx')
        
        result = pdf_to_ppt(input_file, output_file)
        print(result)
    else:
        print("用法: python pdf_to_ppt.py <input.pdf> [output.pptx]")
